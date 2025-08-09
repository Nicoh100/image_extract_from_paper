from __future__ import annotations
import os
import re
import sys
from pathlib import Path
from urllib.parse import urljoin

import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt

BMJ_URL = "https://www.bmj.com/content/388/bmj-2024-080881.long"

CAPTIONS = {
    "Fig 1": "Consolidated standards of reporting trials (CONSORT) diagram. Assessment conducted refers to clinical review, not assessment of primary endpoint",
    "Fig 2": "Amputation free survival (AFS) Kaplan-Meier plot and hazard ratio over time fitted assuming non-proportional hazards (intention-to-treat analysis). CI=confidence interval",
}


def fetch_html(url: str) -> BeautifulSoup:
    headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}
    r = requests.get(url, headers=headers, timeout=30)
    r.raise_for_status()
    return BeautifulSoup(r.text, "html.parser")


def _pick_best_img_src(img) -> str | None:
    # Prefer high-res attributes
    for key in ("data-zoom-src", "data-original"):
        if img.get(key):
            return img.get(key)
    # Parse srcset/data-srcset for highest width
    srcset = img.get("srcset") or img.get("data-srcset")
    if srcset:
        best_url = None
        best_w = -1
        for part in srcset.split(","):
            part = part.strip()
            if not part:
                continue
            bits = part.split()
            url = bits[0]
            w = 0
            if len(bits) > 1 and bits[1].lower().endswith("w"):
                try:
                    w = int(bits[1][:-1])
                except Exception:
                    w = 0
            if w > best_w:
                best_w = w
                best_url = url
        if best_url:
            return best_url
    # Fallback to src/data-src
    return img.get("src") or img.get("data-src")


def find_figure_images(soup: BeautifulSoup):
    results = []
    for fig in soup.find_all(["figure", "div"], class_=re.compile(r"fig|figure", re.I)):
        label = None
        label_el = fig.find(["span", "div", "p", "strong"], class_=re.compile(r"fig|label", re.I))
        if label_el:
            txt = label_el.get_text(" ", strip=True)
            m = re.search(r"Fig(?:ure)?\s*(\d+)", txt, re.I)
            if m:
                label = f"Fig {m.group(1)}"
        caption_el = fig.find(["figcaption", "caption", "div", "p"], class_=re.compile(r"caption", re.I))
        caption_text = caption_el.get_text(" ", strip=True) if caption_el else None
        img = fig.find("img")
        if img:
            src = _pick_best_img_src(img)
            if src:
                results.append({"label": label, "caption": caption_text, "src": src})
    return results


def download_image(base_url: str, src: str, out_dir: Path, name_hint: str) -> Path:
    url = src if src.startswith("http") else urljoin(base_url, src)
    headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}
    r = requests.get(url, headers=headers, timeout=60)
    r.raise_for_status()
    out_dir.mkdir(parents=True, exist_ok=True)
    # Guess extension
    ext = ".png"
    if ".jpg" in url.lower() or ".jpeg" in url.lower():
        ext = ".jpg"
    elif ".svg" in url.lower():
        ext = ".svg"
    fname = re.sub(r"[^A-Za-z0-9_\-]", "_", name_hint) + ext
    out_path = out_dir / fname
    with open(out_path, "wb") as f:
        f.write(r.content)
    return out_path


def build_ppt(images_with_captions: list[tuple[Path, str]], output_path: Path) -> None:
    prs = Presentation()
    slide_w_in = prs.slide_width / 914400.0
    for img_path, caption in images_with_captions:
        slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title_h_in = 0.6
        tx = slide.shapes.add_textbox(Inches(0), Inches(0.1), Inches(slide_w_in), Inches(title_h_in))
        p = tx.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(20)
        from PIL import Image as PILImage
        if img_path.suffix.lower() == ".svg":
            # pptx doesn't support SVG directly; rely on size placeholder
            # Simple fallback: place without resizing
            final_w = slide_w_in - 1.0
            left = 0.5
            top = title_h_in + 0.2
            slide.shapes.add_picture(str(img_path), Inches(left), Inches(top), width=Inches(final_w))
            continue
        with PILImage.open(img_path) as im:
            w, h = im.size
        slide_h_in = prs.slide_height / 914400.0
        usable_w = slide_w_in - 1.0
        usable_h = slide_h_in - (title_h_in + 0.5)
        aspect = w / max(1, h)
        if usable_w / aspect > usable_h:
            final_h = usable_h; final_w = final_h * aspect
        else:
            final_w = usable_w; final_h = final_w / aspect
        left = (slide_w_in - final_w) / 2.0
        top = title_h_in + (usable_h - final_h) / 2.0
        slide.shapes.add_picture(str(img_path), Inches(left), Inches(top), width=Inches(final_w))
    prs.save(str(output_path))


def main():
    out_root = Path("bmj_figs")
    soup = fetch_html(BMJ_URL)
    figs = find_figure_images(soup)
    wanted = {"Fig 1", "Fig 2"}
    # Build mapping label->src (unique), keeping first occurrence
    label_to_src: dict[str, str] = {}
    for f in figs:
        label = f.get("label")
        src = f.get("src")
        if not src or not label:
            continue
        if label in wanted and label not in label_to_src:
            label_to_src[label] = src
        if len(label_to_src) == len(wanted):
            break
    # Fallback: search any image with alt mentioning Fig 1 / Fig 2
    if len(label_to_src) < len(wanted):
        for img in soup.select("img"):
            alt = (img.get("alt") or "").strip()
            src = _pick_best_img_src(img)
            if not src:
                continue
            if "Fig 1" in alt and "Fig 1" not in label_to_src:
                label_to_src["Fig 1"] = src
            elif "Fig 2" in alt and "Fig 2" not in label_to_src:
                label_to_src["Fig 2"] = src
            if len(label_to_src) == len(wanted):
                break
    if not label_to_src:
        print("No target figures found.")
        sys.exit(1)
    # Download and build PPT for exactly Fig 1 and Fig 2 if available
    images_with_caps = []
    for label in sorted(label_to_src.keys(), key=lambda s: int(re.search(r"\d+", s).group(0))):
        src = label_to_src[label]
        out_path = download_image(BMJ_URL, src, out_root, label)
        caption = f"{label} {CAPTIONS.get(label, label)}"
        images_with_caps.append((out_path, caption))
        print(f"Selected {label}: {src}")
    pptx_out = Path("bmj_figs.pptx")
    build_ppt(images_with_caps, pptx_out)
    print(f"Saved {pptx_out} with {len(images_with_caps)} slides.")

if __name__ == "__main__":
    main()
