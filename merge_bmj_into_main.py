from __future__ import annotations
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

CAPTIONS = {
    "Fig 1": "Consolidated standards of reporting trials (CONSORT) diagram. Assessment conducted refers to clinical review, not assessment of primary endpoint",
    "Fig 2": "Amputation free survival (AFS) Kaplan-Meier plot and hazard ratio over time fitted assuming non-proportional hazards (intention-to-treat analysis). CI=confidence interval",
}


def add_slide_with_image(prs: Presentation, image_path: Path, caption: str) -> None:
    slide_w_in = prs.slide_width / 914400.0
    slide_h_in = prs.slide_height / 914400.0

    # margins similar to earlier logic (12,12,14,14 mm)
    left_in, top_in, right_in, bottom_in = (12/25.4, 12/25.4, 14/25.4, 14/25.4)
    usable_w = slide_w_in - (left_in + right_in)
    usable_h = slide_h_in - (top_in + bottom_in)

    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    title_h = 0.6
    tx = slide.shapes.add_textbox(Inches(0), Inches(max(0.1, top_in - 0.2)), Inches(slide_w_in), Inches(title_h))
    p = tx.text_frame.paragraphs[0]
    p.text = caption
    p.font.size = Pt(20)

    with Image.open(image_path) as im:
        w, h = im.size
    aspect = w / max(1, h)
    max_w = usable_w
    max_h = usable_h - title_h
    if max_h <= 0:
        max_h = usable_h

    if max_w / aspect > max_h:
        final_h = max_h
        final_w = final_h * aspect
    else:
        final_w = max_w
        final_h = final_w / aspect

    left = left_in + max(0.0, (usable_w - final_w) / 2.0)
    top = top_in + title_h + max(0.0, (max_h - final_h) / 2.0)

    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(final_w))


def main():
    # Inputs
    deck_path = Path(r"C:\Users\nico-\Desktop\Maren Präsentation\py script\_Plain versus drug balloon and stenting in severe ischaemia of the - Kopie_cropped_extracted.pptx")
    bmj_dir = Path(r"C:\Users\nico-\Desktop\Maren Präsentation\py script\bmj_figs")
    fig1 = bmj_dir / "Fig_1.jpg"
    fig2 = bmj_dir / "Fig_2.jpg"

    if not deck_path.exists():
        raise SystemExit(f"Deck not found: {deck_path}")
    if not fig1.exists() or not fig2.exists():
        raise SystemExit("BMJ figure images not found. Run bmj_fig_downloader.py first.")

    prs = Presentation(str(deck_path))
    add_slide_with_image(prs, fig1, f"Fig 1 {CAPTIONS['Fig 1']}")
    add_slide_with_image(prs, fig2, f"Fig 2 {CAPTIONS['Fig 2']}")

    prs.save(str(deck_path))
    print(f"Merged BMJ figures into: {deck_path}")


if __name__ == "__main__":
    main()
