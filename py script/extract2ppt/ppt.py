from __future__ import annotations

from pathlib import Path
from typing import List, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from rich.console import Console

from .types import AssetMeta, PipelineConfig

console = Console()


def _slide_size_inches(prs: Presentation) -> Tuple[float, float]:
    return prs.slide_width / 914400.0, prs.slide_height / 914400.0


def _add_title(slide, text: str, slide_width_in: float, top_margin_in: float) -> float:
    title_height_in = 0.6
    left = Inches(0)
    top = Inches(top_margin_in - 0.2) if top_margin_in > 0.25 else Inches(0.1)
    width = Inches(slide_width_in)
    height = Inches(title_height_in)
    tx = slide.shapes.add_textbox(left, top, width, height)
    p = tx.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(20)
    return title_height_in


def build_ppt(
    assets: List[AssetMeta],
    config: PipelineConfig,
    exports_root: Path,
    output_pptx: Path,
) -> None:
    prs = Presentation(config.template_path) if config.template_path else Presentation()
    slide_w_in, slide_h_in = _slide_size_inches(prs)

    left_mm, top_mm, right_mm, bottom_mm = config.margins_mm
    left_in = left_mm / 25.4
    top_in = top_mm / 25.4
    right_in = right_mm / 25.4
    bottom_in = bottom_mm / 25.4

    usable_w_in = slide_w_in - (left_in + right_in)
    usable_h_in = slide_h_in - (top_in + bottom_in)

    for meta in assets:
        slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        title_text = meta.caption or f"{meta.kind.capitalize()} p{meta.page+1}"
        title_h_in = _add_title(slide, title_text, slide_w_in, top_in)

        max_w_in = usable_w_in
        max_h_in = usable_h_in - title_h_in
        if max_h_in <= 0:
            max_h_in = usable_h_in

        img_w_px, img_h_px = meta.width_px, meta.height_px
        img_aspect = img_w_px / max(1, img_h_px)
        if max_w_in / img_aspect > max_h_in:
            final_h_in = max_h_in
            final_w_in = final_h_in * img_aspect
        else:
            final_w_in = max_w_in
            final_h_in = final_w_in / img_aspect

        effective_ppi = img_w_px / max(0.01, final_w_in)
        if effective_ppi < config.min_ppi:
            console.print(f"[yellow]Warning:[/] Asset {meta.filename} effective PPI {effective_ppi:.0f} < {config.min_ppi}")

        left_inch = left_in + max(0.0, (usable_w_in - final_w_in) / 2.0)
        top_inch = top_in + title_h_in + max(0.0, (max_h_in - final_h_in) / 2.0)

        kind_dir = {
            "table": exports_root / "tables",
            "image": exports_root / "images",
            "figure": exports_root / "figures",
        }[meta.kind]
        img_path = kind_dir / meta.filename

        slide.shapes.add_picture(
            str(img_path),
            Inches(left_inch),
            Inches(top_inch),
            width=Inches(final_w_in),
            height=None,
        )

    prs.save(str(output_pptx))
    console.print(f"[green]Saved PPTX:[/] {output_pptx}")
